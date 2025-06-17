from fastapi import FastAPI, UploadFile, File, HTTPException, Form
from fastapi.responses import JSONResponse
import aiofiles
import os
from pptx import Presentation
import pytesseract
from pdf2image import convert_from_path, convert_from_bytes
from openai import OpenAI
from pydantic import BaseModel
import re
# from sentence_transformers import SentenceTransformer, util
from konlpy.tag import Okt
from typing import Dict
from pdf2image import convert_from_path
import requests
from fastapi.middleware.cors import CORSMiddleware


#pdf2image 설정
poppler_path=r"C:/Users/jhs38/poppler-24.08.0/Library/bin"

# KorSBERT 모델, Okt 형태소 분석기
# model = SentenceTransformer('snunlp/KR-SBERT-V40K-klueNLI-augSTS')
#okt = Okt()

#py -m pip install fastapi uvicorn aiofiles pdf2image pytesseract python-pptx openai konlpy
#py -m uvicorn main:app --reload --port 8000

# openai api key 설정
client = OpenAI(api_key="") #api key를 여기에 입력하세요

# 업로드 디렉토리 설정
UPLOAD_DIR = "uploads"
os.makedirs(UPLOAD_DIR, exist_ok=True)

#OCR용 tesseract 경로 설정
pytesseract.pytesseract.tesseract_cmd = r"C:/Program Files/Tesseract-OCR/tesseract.exe"
    
# # Spring STT(Whisper) API 엔드포인트
# SPRING_STT_URL = "http://localhost:8080/api/whisper-multi" # Spring Boot 서버의 STT API 엔드포인트

# 의도 일치 임계값
INTENT_THRESHOLD = 0.75


# stt 비교 함수 ---------------------------------------------------------------------------------------------
class TextRequest(BaseModel):
    raw_text: str
    
# ✅ 입력 데이터 모델
class CompareRequest(BaseModel):
    cue: str      # 큐시트 문장
    stt: str      # STT 결과 문장

# ✅ 핵심 키워드 추출 함수
# def get_keywords(text: str) -> set:
#     return set(okt.nouns(text))

class EvaluateAudioResponse(BaseModel):
    stt_text: str
    sentence_similarity: float
    intent_match: float
    keyword_coverage: float
    level: str       # “높음”/“중간”/“낮음” 등급


# 텍스트 품질 검사 함수 --------------------------------------------------------------------------------------------


# 영어 단어/불용어 관련 코드는 사용하지 않음
def clean_text(text):
    return re.sub(r'[^가-힣a-zA-Z0-9\s]', '', text)

def filter_english_words(text):
    return re.findall(r'\b[a-zA-Z]{2,}\b', text)

def filter_korean_words(text):
    nouns = okt.nouns(text)
    return [n for n in nouns if len(n) > 1]

def extract_valid_words(text):
    text = clean_text(text)
    eng_words = filter_english_words(text)
    kor_words = filter_korean_words(text)
    return eng_words, kor_words
    
    
#텍스트 추출 함수 ----------------------------------------------------------------------------------------------

#pdf
def extract_text_from_pdf(pdf_input, from_bytes: bool = False) -> str:
    if from_bytes:
        images = convert_from_bytes(pdf_input)
    else:
        images = convert_from_path(pdf_input)
    full_text = ""
    for img in images:
        text = pytesseract.image_to_string(img, lang='kor+eng')
        full_text += text + "\n"
    return full_text


#pptx 
def extract_text_from_pptx(path: str) -> str:
    text = ""
    prs = Presentation(path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
            elif hasattr(shape, "text_frame") and shape.text_frame is not None:
                text += shape.text_frame.text + "\n"
    return text

# 문장 유사도에 따라 "높음"/"중간"/"낮음" 등급을 결정
def determine_level(similarity: float) -> str:
    if similarity >= 0.75:
        return "높음"
    elif similarity >= 0.50:
        return "중간"
    else:
        return "낮음"


# 스크립트 --------------------------------------------------------------------------------------------
def generate_presentation_script(slide_text: str) -> str:
    prompt = f"""
너는 대학생 발표자료에서 발표자가 사용할 발표 대본과 요약 큐카드를 생성하는 전문가야.

아래 슬라이드별 텍스트를 참고해, 반드시 아래 지침과 예시를 따라 발표자료 대본을 작성해 줘.

---
[지침]
[0] 슬라이드 번호 자동 지정
- 입력된 텍스트에서 슬라이드별 명확한 구분자(예: 제목) 없이도, 각 주요 섹션 순서대로 슬라이드 번호(1, 2, 3, ...)를 부여하여 모든 슬라이드를 누락 없이 처리하세요.

[1] 기본버전
- 슬라이드의 텍스트와 구조를 참고하여, 각 페이지(슬라이드)마다 발표자가 읽을 수 있는 발표 대본을 작성하세요.
- 문장은 자연스럽고 논리적인 흐름을 갖추되, 단정적이고 공식적인 발표체(예: “~입니다”, “~합니다”)를 사용하세요.
- 구어체(예: “~거든요”, “~해요”, “~같습니다”)는 사용하지 마세요.
- 슬라이드 제목만을 주제로 삼거나, 파일명으로 내용을 유추하지 마세요.
- 발표자가 실제로 말하지 않을 내용(예: “이 슬라이드는 ~를 보여줍니다” 등 슬라이드 설명이나 화면 안내)은 포함하지 마세요.
- 표지(1번 슬라이드)에서는 인사와 발표 주제를 간단하게 안내하는 수준으로만 작성하세요.
- OCR 인식이 불가능하거나 텍스트가 거의 없는 슬라이드는 아래와 같이 작성하세요:
    [기본버전] (OCR 인식 불가 - 요약 생략)
- 발표 흐름이 자연스럽게 이어질 수 있도록, 각 슬라이드 마지막 문장은 다음 내용을 예고하거나, 적절한 연결 어미로 마무리하세요.
- 발표자가 활용할 수 있도록 아래 비언어적 표현 아이콘을 적절한 위치(강조, 전환, 호흡 등)에 배치하세요:
    🔍 청중 바라보기
    📄 발표자료 보기
    ✋ 제스처
    👉 화면 가리키기
    🌬 호흡
    ❓ 질의응답
- 아이콘은 실제 발표 흐름에 어울리는 위치에 자연스럽게 삽입하세요.
    예시: 슬라이드 첫 문장 앞(호흡/청중 바라보기), 중요한 정보 뒤(화면 가리키기/제스처), 주제 전환 시(호흡/자료 보기) 등

[2] 심화버전
- 해당 슬라이드의 핵심 주제를 한문장으로 논리적으로 요약하세요.
- 요약문은 간단한 문장으로, **핵심 키워드만** 포함하여 작성하세요.
- 주요 메시지를 빠르게 파악할 수 있도록, 선언문 또는 설명문 형태로 간결하게 작성하세요.
- 불필요한 부연설명은 피하고, 문서의 핵심 논지에 집중하세요.
- 표지(1번 슬라이드)는 ‘이번 발표의 목적’만 간결하게 설명하세요.
- OCR 인식이 불가능한 경우 아래처럼 작성하세요:
    [심화버전] (OCR 인식 불가 – 요약 생략)

[3] 출력 형식 예시

슬라이드 1
[기본버전]
안녕하세요. 오늘 ‘사업 타당성 분석’에 대해 발표할 경영학과 20210001 김지원입니다. 🌬 호흡
지금부터 발표를 시작하겠습니다. 🔍 청중 바라보기

[심화버전]
발표 주제 설명 및 자기소개

슬라이드 2
[기본버전]
타당성 분석은 사업 아이디어가 실제로 실행 가능한지 판단하는 절차입니다.
이를 통해 사업화 가치가 있는지 평가할 수 있습니다. 👉 화면 가리키기
이제 타당성 분석의 시기와 중요성에 대해 설명드리겠습니다. 🌬 호흡

[심화버전]
타당성분석 개념 및 목적 설명

슬라이드 3
[기본버전]
타당성 분석은 사업 아이디어의 실행 가능성을 예비 평가하는 것으로, 비즈니스의 초기 단계에서 수행해야 가장 효과적입니다.
많은 리소스가 투입되기 전에 아이디어를 선별하는 데 중요한 역할을 합니다. 🔍 청중 바라보기
일부 기업가는 아이디어 파악 후 바로 비즈니스 모델 개발로 넘어가 실수를 할 수 있습니다. ✋ 제스처
효과적인 타당성 분석은 이러한 오류를 방지할 수 있습니다. 🌬 호흡

[심화버전]
타다성 분석을 비즈니스 초기단계 수행해야하는 중요성을 설명

[4] 추가 지침
- 각 슬라이드별로 [기본버전], [심화버전] 순서로 반드시 출력하세요.
- 형식, 아이콘, 슬라이드/버전 순서를 꼭 지켜주세요.
- 출력 결과가 위 예시와 다르거나 형식이 어긋나지 않도록 주의하세요.

---
[슬라이드별 OCR 또는 텍스트 추출 결과]
{slide_text}

---
[최종 출력]
"""
    response = client.chat.completions.create(
        model="gpt-4.1-nano",
        messages=[
            {
                "role": "system",
                "content": "너는 실전 발표 큐카드/대본 자동 생성 전문가야. 반드시 예시 포맷과 지침을 지켜서, 발표자가 발표자료만 보고도 완성도 높은 발표를 할 수 있게 작성해."
            },
            {
                "role": "user",
                "content": prompt
            }
        ],
        temperature=0.1
    )
    return response.choices[0].message.content.strip()

#───(GPT에게 직접 비교 요청하는 헬퍼 함수)────────────────────────────────────

def gpt_compare(cue: str, stt: str) -> Dict:
    """
    GPT에게 'cue'와 'stt' 두 문장을 주고,
    JSON 형태(키: sentence_similarity, intent_match, keyword_coverage, level)로
    결과만 딱 리턴해 달라고 요청합니다.
    """
    prompt = f"""
아래 두 문장(기준 문장, STT 결과 문장)에 대해 다음 네 가지 지표를 계산하고, JSON 객체 형태로 출력해주세요.

1) sentence_similarity: 0.00에서 1.00 사이의 숫자로, 두 문장의 의미적 유사도를 나타내세요.
   - 1.00: 의미가 거의 동일
   - 0.00: 전혀 다른 의미

2) intent_match: 0.00에서 1.00 사이의 숫자로, 두 문장의 의도가 얼마나 일치하는지 나타내세요.
   - 예: 두 문장의 핵심 의도가 같다면 1.00, 완전히 다르면 0.00

3) keyword_coverage: 0.00에서 1.00 사이의 숫자로, 기준 문장(cue)에 있는 핵심 키워드가 STT 문장에 얼마나 포함되었는지 비율로 나타내세요.
   - 핵심 키워드는 GPT가 문맥을 보고 추출하도록 하세요.

4) level: 다음 기준에 따라 "높음" / "중간" / "낮음" 중 하나를 출력하세요.
   - sentence_similarity가 0.75 이상이면 "높음"
   - 0.50 이상 0.75 미만이면 "중간"
   - 그 외에는 "낮음"

반환 예시(JSON만 출력; 추가 설명이나 텍스트 없이 딱 JSON):

[기준 문장]
"{cue}"

[STT 결과 문장]
"{stt}"
"""
    response = client.chat.completions.create(
        model="gpt-4.1-nano",
        messages=[
            {"role": "user", "content": prompt}
        ],
        temperature=0.0
    )
    content = response.choices[0].message.content.strip()

    # GPT가 출력한 JSON만 파싱
    try:
        import json
        return json.loads(content)
    except Exception:
        # JSON 파싱 오류 시, 매우 간단히 숫자와 레벨을 추출
        # (실 시연에서는 JSON 형식 잘 지켜달라고 Prompts에 강조하면 됨)
        m = {}
        nums = re.findall(r'"sentence_similarity"\s*:\s*([0-9.]+)', content)
        m["sentence_similarity"] = float(nums[0]) if nums else 0.0
        
        # nums = re.findall(r'"intent_match"\s*:\s*([0-9.]+)', content)
        # m["intent_match"] = float(nums[0]) if nums else 0.0
        # nums = re.findall(r'"keyword_coverage"\s*:\s*([0-9.]+)', content)
        # m["keyword_coverage"] = float(nums[0]) if nums else 0.0
        
        lvl = re.search(r'"level"\s*:\s*"([^"]+)"', content)
        m["level"] = lvl.group(1) if lvl else ""
        return m


# ───(FastAPI 앱 정의)────────────────────────────────────────────

app = FastAPI(
    title="통합 OCR·큐카드 생성·STT 비교 API",
    description="PDF/PPTX에서 텍스트 추출, 발표 스크립트 생성, STT vs 큐시트 비교를 한 곳에서 처리합니다.",
    version="1.0.0"
)

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],    # 나중에 운영 환경에 맞게 도메인 지정
    allow_methods=["*"],
    allow_headers=["*"],
)

@app.get("/")
def root():
    return {"message": "API is running"}


@app.post("/extract_words/pdf")
async def extract_words_from_pdf(file: UploadFile = File(...)):
    """
    PDF → OCR(text) → 영어/한국어 단어 집합 반환
    """
    contents = await file.read()
    full_text = extract_text_from_pdf(contents, from_bytes=True)
    eng_words, kor_words = extract_valid_words(full_text)
    return JSONResponse(content={
        "english_words": eng_words,
        "korean_words": kor_words
    })


@app.post("/extract-text-pdf-pptx/")
async def extract_ocr_text(file: UploadFile = File(...)):
    """
    업로드된 PDF/PPTX → OCR/text 추출 후 반환
    """
    filename = file.filename
    ext = filename.rsplit(".", 1)[-1].lower()
    saved_path = os.path.join(UPLOAD_DIR, filename)

    async with aiofiles.open(saved_path, "wb") as out_file:
        await out_file.write(await file.read())

    try:
        if ext == "pdf":
            text = extract_text_from_pdf(saved_path, from_bytes=False)
        elif ext == "pptx":
            text = extract_text_from_pptx(saved_path)
        else:
            raise HTTPException(status_code=400, detail="PDF 또는 PPTX만 지원합니다.")
        return JSONResponse({
            "filename": filename,
            "text_by_ocr": text
        })
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"처리 중 오류 발생: {str(e)}")


@app.post("/generate-presentation-script/")
async def generate_presentation_script_api(req: TextRequest):
    """
    슬라이드 텍스트(raw_text) → 발표 대본 & 큐카드 자동 생성
    """
    raw = req.raw_text
    if not raw.strip():
        raise HTTPException(status_code=400, detail="raw_text가 비어 있습니다.")

    try:
        result = generate_presentation_script(raw)
        return JSONResponse({
            "presentation_script": result
        })
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"처리 중 오류 발생: {str(e)}")


# @app.post("/evaluate", summary="큐시트 vs STT 비교 (문자열 입력)", response_model=Dict[str, float])
# def evaluate_similarity(data: CompareRequest):
#     """
#     요청 JSON: { "cue": "...", "stt": "..." }
#     - 1) KorSBERT 코사인 유사도
#     - 2) Threshold(0.75) 기반 의도 일치 (1.0/0.0)
#     - 3) 키워드 포함률(자카드 유사도)
#     """
#     cue = data.cue
#     stt = data.stt

#     # 1. KorSBERT로 임베딩 → 코사인 유사도 계산
#     emb = model.encode([cue, stt])
#     similarity = float(util.cos_sim(emb[0], emb[1])[0][0])

#     # 2. 의도 일치(Threshold)
#     intent_match = 1.0 if similarity >= INTENT_THRESHOLD else 0.0

#     # 3. 키워드 포함률 (자카드 유사도)
#     cue_keywords = get_keywords(cue)
#     stt_keywords = get_keywords(stt)
#     intersection = cue_keywords & stt_keywords
#     union = cue_keywords | stt_keywords
#     keyword_coverage = len(intersection) / len(union) if union else 0.0

#     return {
#         "sentence_similarity": round(similarity, 4),
#         "intent_match": intent_match,
#         "keyword_coverage": round(keyword_coverage, 4)
#     }


# @app.post(
#     "/evaluate-with-audio",
#     summary="큐시트 문장 + 오디오 파일(STT) → GPT 비교",
#     response_model=EvaluateAudioResponse
# )
# async def evaluate_with_audio(
#     cue: str = Form(..., description="비교할 기준(큐시트) 문장"),
#     audio: UploadFile = File(..., description="STT 대상 오디오 파일 (예: WAV)")
# ):
#     # 1) 오디오 파일 임시 저장
#     temp_path = os.path.join(UPLOAD_DIR, f"temp_{audio.filename}")
#     async with aiofiles.open(temp_path, "wb") as out_file:
#         await out_file.write(await audio.read())

#     # 2) Spring STT API 호출해서 stt_text 받아오기
#     try:
#         with open(temp_path, "rb") as f:
#             files = {"audio": f}
#             resp = requests.post(SPRING_STT_URL, files=files, timeout=60)
#             resp.raise_for_status()
#     except Exception as e:
#         if os.path.exists(temp_path):
#             os.remove(temp_path)
#         raise HTTPException(status_code=500, detail=f"STT API 호출 중 오류 발생: {e}")

#     data = resp.json()
#     stt_text = data.get("result", "").strip()

#     # 3) 임시 파일 삭제
#     if os.path.exists(temp_path):
#         os.remove(temp_path)

#     # 4) GPT에게 비교 요청
#     r = gpt_compare(cue, stt_text)

#     # 5) GPT가 만들어준 결과를 EvaluateAudioResponse 형태로 반환
#     return EvaluateAudioResponse(
#         stt_text=stt_text,
#         sentence_similarity=r.get("sentence_similarity", 0.0),
#         intent_match=r.get("intent_match", 0.0),
#         keyword_coverage=r.get("keyword_coverage", 0.0),
#         level=r.get("level", "")
#     )

@app.post(
    "/evaluate-with-audio",
    summary="큐시트 문장 + 오디오 파일(STT) → GPT 비교",
    response_model=EvaluateAudioResponse
)
async def evaluate_with_audio(
    cue: str = Form(..., description="비교할 기준(큐시트) 문장"),
    audio: UploadFile = File(..., description="STT 대상 오디오 파일 (예: WAV)")
):
    # 오디오 파일 임시 저장
    temp_path = os.path.join(UPLOAD_DIR, f"temp_{audio.filename}")
    async with aiofiles.open(temp_path, "wb") as out_file:
        await out_file.write(await audio.read())

    # 여기부터 예외가 발생하면 에러 메시지를 바로 반환하도록 수정
    try:
        # 1) Whisper API 호출 부분
        with open(temp_path, "rb") as fbin:
            whisper_resp = client.audio.transcriptions.create(
                file=fbin,
                model="whisper-1"
            )
        # Whisper 응답에서 텍스트 꺼내기
        stt_text = whisper_resp.text.strip()

        # 2) 임시 파일 삭제
        if os.path.exists(temp_path):
            os.remove(temp_path)

        # 3) GPT 비교
        r = gpt_compare(cue, stt_text)

        return EvaluateAudioResponse(
            stt_text=stt_text,
            sentence_similarity=r.get("sentence_similarity", 0.0),
            #intent_match=r.get("intent_match", 0.0),
            #keyword_coverage=r.get("keyword_coverage", 0.0),
            level=r.get("level", "")
        )

    except Exception as e:
        # 예외 메시지를 자세히 남기고 500으로 응답
        # 실제 시연이 끝나면 e를 그대로 보여주는 코드는 삭제하세요.
        detail_msg = f"예외 발생: {type(e).__name__} – {e}"
        # 임시 파일도 남지 않도록 삭제
        if os.path.exists(temp_path):
            os.remove(temp_path)
        raise HTTPException(status_code=500, detail=detail_msg)

# 시연용 test API

# 시연용으로 미리 만들어 둔 카드 데이터
DEMO_CARDS = [
    { "title": "슬라이드 1", 
     "content":  " \n[기본버전]  \n\n안녕하세요. 오늘은 대학생 발표 불안과 그 해결 방안에 대해 발표하겠습니다. 🌬 호흡  \n발표 주제는 대학생들이 발표를 준비하고 수행하는 과정에서 겪는 어려움과 이를 극복하는 방법입니다. 🔍 청중 바라보기\n\n\n \n발표 불안은 발표할 내용을 기억하지 못하거나, 말이 떠오르지 않거나, 사람들 앞에서 말하는 것을 피하게 되는 두려움으로 나타납니다. ✋ 제스처  \n이러한 두려움은 발표 연단에서의 자신감 저하와 발표 실패로 이어질 수 있습니다. 🌬 호흡  \n이 문제를 해결하기 위해서는 적절한 준비와 연습이 필요합니다. 📄 발표자료 보기\n\n  \n  \n대학 신입생 117명을 대상으로 한 설문조사 결과, 69.6%가 발표 수업이 부담스럽다고 응답했습니다. 🔍 청중 바라보기  \n또한, 62.4%는 모든 문장을 미리 스크립트로 작성한다고 답변했으며, 이는 발표 준비의 한 방법입니다. 🌬 호흡  \n이러한 조사 결과는 발표 불안을 줄이기 위한 전략 수립에 중요한 자료가 됩니다. 📄 발표자료 보기\n\n" }]

@app.post("/api/qcards")
async def demo_qcards():
    return JSONResponse(content={ "data": DEMO_CARDS })


# 요청 바디 스키마
class EvaluateTextRequest(BaseModel):
    cue: str
    transcripts: list[str]

# 응답 항목 스키마
class EvaluateResponse(BaseModel):
    sentence_similarity: float
    level: str
    # 필요하면 intent_match, keyword_coverage 등 추가

@app.post("/api/evaluate", response_model=list[EvaluateResponse])
async def evaluate_text(req: EvaluateTextRequest):
    try:
        # cue, transcripts 받음
        cue = req.cue
        transcripts = req.transcripts

        # gpt_compare 는 cue와 단일 문장을 받아서
        # {"sentence_similarity":..., "level":...} 반환하는 헬퍼 함수
        results = [gpt_compare(cue, stt) for stt in transcripts]

        return results

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))
